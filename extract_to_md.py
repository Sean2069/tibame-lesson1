"""
從 files/ 資料夾讀取各種格式文件，萃取文字後以 Markdown 格式儲存至 doc/ 資料夾。
支援格式：PDF、PPTX、DOCX、MD
"""

import os
import sys
import re
from pathlib import Path

import pdfplumber
from pptx import Presentation
from docx import Document


FILES_DIR = Path(__file__).parent / "files"
DOC_DIR = Path(__file__).parent / "doc"


# ── 各格式萃取函式 ──────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    lines = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                lines.append(f"<!-- page {i} -->")
                lines.append(text.strip())
    return "\n\n".join(lines)


def extract_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    parts.append(text)
        if parts:
            slides.append((i, parts))
    return slides  # 回傳結構化資料，由 to_markdown 處理


def extract_docx(path: Path) -> str:
    doc = Document(path)
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paras)


def extract_md(path: Path) -> str:
    return path.read_text(encoding="utf-8")


# ── Markdown 產生函式 ────────────────────────────────────────────────────────

def guess_heading_level(line: str, prev_lines: list[str]) -> str:
    """根據簡單啟發式判斷一行是否為標題，回傳 Markdown 標題格式或原文。"""
    stripped = line.strip()
    if not stripped:
        return ""

    # 已有 Markdown 標題符號 → 保留
    if stripped.startswith("#"):
        return stripped

    # 全大寫、短句（< 30 字元）且不以標點結尾 → 視為標題
    is_short = len(stripped) <= 30
    ends_with_punct = stripped[-1] in "。，、：:,.!?！？"
    has_digit_prefix = re.match(r"^[\d一二三四五六七八九十]+[、.]", stripped)

    if is_short and not ends_with_punct:
        return f"## {stripped}"
    if has_digit_prefix:
        return f"### {stripped}"
    return stripped


def pdf_to_markdown(raw: str, title: str) -> str:
    lines = raw.split("\n")
    md_lines = [f"# {title}", ""]
    prev = []
    for line in lines:
        if line.startswith("<!-- page"):
            page_num = re.search(r"\d+", line).group()
            md_lines.append(f"\n---\n\n> 第 {page_num} 頁\n")
            continue
        result = guess_heading_level(line, prev)
        md_lines.append(result)
        prev.append(line)
    return "\n".join(md_lines)


def pptx_to_markdown(slides_data: list, title: str) -> str:
    md_lines = [f"# {title}", ""]
    for slide_num, parts in slides_data:
        md_lines.append(f"\n## 第 {slide_num} 張投影片\n")
        if parts:
            # 第一個文字方塊通常是投影片標題
            md_lines.append(f"### {parts[0]}\n")
            for part in parts[1:]:
                # 多行內容轉為清單
                for sub in part.split("\n"):
                    sub = sub.strip()
                    if sub:
                        if len(sub) <= 30 and not sub[-1] in "。，、：:,.!?！？":
                            md_lines.append(f"#### {sub}")
                        else:
                            md_lines.append(f"- {sub}")
    return "\n".join(md_lines)


def docx_to_markdown(raw: str, title: str) -> str:
    lines = raw.split("\n")
    md_lines = [f"# {title}", ""]
    prev = []
    for line in lines:
        result = guess_heading_level(line, prev)
        md_lines.append(result if result else "")
        prev.append(line)
    return "\n".join(md_lines)


def md_to_markdown(raw: str, title: str) -> str:
    # 若原始 MD 已有一級標題則保留，否則插入
    if raw.lstrip().startswith("# "):
        return raw
    return f"# {title}\n\n{raw}"


# ── 主流程 ──────────────────────────────────────────────────────────────────

def process_file(src: Path, dst_dir: Path):
    ext = src.suffix.lower()
    stem = src.stem  # 不含副檔名的檔名，作為標題

    print(f"  處理：{src.name} ({ext})")

    if ext == ".pdf":
        raw = extract_pdf(src)
        md = pdf_to_markdown(raw, stem)
    elif ext == ".pptx":
        slides_data = extract_pptx(src)
        md = pptx_to_markdown(slides_data, stem)
    elif ext == ".docx":
        raw = extract_docx(src)
        md = docx_to_markdown(raw, stem)
    elif ext == ".md":
        raw = extract_md(src)
        md = md_to_markdown(raw, stem)
    else:
        print(f"    略過不支援的格式：{ext}")
        return

    out_path = dst_dir / f"{stem}.md"
    out_path.write_text(md, encoding="utf-8")
    print(f"    → 儲存至：{out_path.relative_to(dst_dir.parent)}")


def main():
    if not FILES_DIR.exists():
        print(f"找不到 files/ 資料夾：{FILES_DIR}")
        sys.exit(1)

    DOC_DIR.mkdir(exist_ok=True)
    print(f"輸出資料夾：{DOC_DIR}\n")

    supported = {".pdf", ".pptx", ".docx", ".md"}
    files = [f for f in sorted(FILES_DIR.iterdir()) if f.suffix.lower() in supported]

    if not files:
        print("files/ 中沒有支援的文件。")
        return

    for f in files:
        process_file(f, DOC_DIR)

    print(f"\n完成！共處理 {len(files)} 個文件。")


if __name__ == "__main__":
    main()
